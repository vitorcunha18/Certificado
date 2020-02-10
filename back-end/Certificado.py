# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

from xlrd import *

from back.BANCO import SGBD


#Lendo planilhas do excel
class Excel(object):
    def __init__(self, nome):
        x = Certificado()
        self.nome = nome
        self.arquivo = open_workbook(self.nome)
        self.planilha = self.arquivo.sheet_by_index(0)
        if linha in range (self.planilha.nrows):
            self.__local = x.run(self.planilha.row_values()[0])
            bd.Insert_certificado(self.planilha.row_values()[0], self.planilha.row_values()[1], self.__local)


#Criação de certificado
class Certificado(object):
    def __init__(self,modelo_power, endereco_salvar, titulo, data, tex, nome_pas=None):
        self.modelo_power, self.local_salvar, self.nome_pasta = modelo_power, endereco_salvar, nome_pas
        self.titulo, self.data, self.tex =  titulo, data, tex,
        if not os.path.exists(f"{self.local_salvar}\\{self.nome_pasta}"):
                    os.mkdir(f"{self.local_salvar}\\{self.nome_pasta}")
        else:
            pass
        self.texto = [f"\n{self.titulo}\n", f'{self.tex}\n', f"{self.data}"]

    def run(self, nome_parti):
        try:
            #Criação e caracterização do caixa de texto 
            self.arq = Presentation(f'{self.modelo_power}')
            self.slide = self.arq.slides[0]
            self.left, self.top, self.width, self.height = Inches(0.5),Inches(1),Inches(9.7),Inches(5)
            self.txBox = self.slide.shapes.add_textbox(self.left, self.top, self.width, self.height)
            self.txBox.margin_top = Inches(5)
            self.tf = self.txBox.text_frame
            self.tf.font_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            self.tf.word_wrap = None

            #adicionando texto e justificando texto
            self.a = self.tf.add_paragraph()
            self.a.alignment = PP_ALIGN.CENTER
            self.a.text = self.texto[0]
            self.a.font.bold = True
            self.a.font.size = Pt(26)
        
            self.a = self.tf.add_paragraph()
            self.a.alignment = PP_ALIGN.JUSTIFY
            self.a.margin_left = Inches(0.5)
            self.a.text = self.texto[1].format(nome_people.upper())
            self.a.font.size = Pt(24)
            self.a = self.tf.add_paragraph()
            self.a.alignment = PP_ALIGN.RIGHT
            self.a.text = self.texto[2]
            self.a.font.size = Pt(24)

            #salvando arquivo
            if "\n" in nome_parti:
                self.nome = str(nome_parti).replace("\n","")
                self.local = f"{self.local_salvar}\\{self.nome_pasta}\\{self.nome}.pptx"
                self.arq.save(self.local)
                return self.local
            else:
                self.local = f"{self.local_salvar}\\{self.nome_pasta}\\{nome_parti}.pptx"
                self.arq.save(self.local)
                return self.local
            


        except:
            print('error Valor')

bd = S
c = Excel('C:\\Users\\Vitor Cunha\\Desktop\\x.xlsx')